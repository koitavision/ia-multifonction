import os
import requests
from flask import Flask, render_template, request, send_file
from dotenv import load_dotenv
from pptx import Presentation
from werkzeug.utils import secure_filename

load_dotenv()

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "static/uploads"

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = os.getenv("OPENROUTER_MODEL", "mistralai/mistral-7b-instruct")

@app.route("/", methods=["GET", "POST"])
def index():
    response_text = ""
    image_url = ""
    ppt_generated = False
    uploaded_file_url = None

    if request.method == "POST":
        task = request.form.get("task")
        prompt = request.form.get("prompt")

        if task == "text":
            response_text = generate_text(prompt)
        elif task == "image":
            image_url = generate_image(prompt)
        elif task == "ppt":
            ppt_path = generate_ppt(prompt)
            ppt_generated = True

        # Gestion de l'upload
        if "file" in request.files:
            file = request.files["file"]
            if file.filename != "":
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
                file.save(filepath)
                uploaded_file_url = "/" + filepath

    return render_template("index.html", response=response_text, image_url=image_url,
                           ppt_generated=ppt_generated, file_url=uploaded_file_url)

@app.route("/download")
def download():
    return send_file("output/generated.pptx", as_attachment=True)

def generate_text(prompt):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }
    json_data = {
        "model": OPENROUTER_MODEL,
        "messages": [{"role": "user", "content": prompt}]
    }
    try:
        response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=json_data)
        data = response.json()
        return data['choices'][0]['message']['content']
    except Exception as e:
        return f"Erreur : {str(e)}"

def generate_image(prompt):
    # Simulation d'image - À remplacer par vraie API image si souhaité
    return f"https://via.placeholder.com/512x512.png?text={prompt.replace(' ', '+')}"

def generate_ppt(prompt):
    prs = Presentation()
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Titre {i + 1}"
        slide.placeholders[1].text = f"Contenu : {prompt} (slide {i + 1})"
    ppt_path = "output/generated.pptx"
    prs.save(ppt_path)
    return ppt_path

if __name__ == "__main__":
    app.run(debug=True)